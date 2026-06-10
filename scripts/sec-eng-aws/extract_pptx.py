#!/usr/bin/env python3
"""
Extract per-slide text, speaker notes, and embedded images from the
Security Engineering on AWS instructor PPTX decks (M00-M09 only; PDFs skipped).

Outputs (to OUT_DIR, kept OUTSIDE the git repo):
  slides.jsonl   - one record per slide: {deck, module, slide_no, title, text, notes}
  images/        - every extracted image, named <module>_s<slide>_i<idx>.<ext>
  images.json    - manifest of every image with slide ref + EMU/inch dimensions
  summary.json   - per-deck counts
"""
import glob, json, os, re, hashlib
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = "/Users/schinchli/Documents/Training Material /Security Engineering on AWS"
OUT = os.path.join(SRC, "_extracted")
IMG_DIR = os.path.join(OUT, "images")
os.makedirs(IMG_DIR, exist_ok=True)

MODULE_RE = re.compile(r"^(M\d{2})_")


def module_id(path):
    m = MODULE_RE.match(os.path.basename(path))
    return m.group(1) if m else os.path.basename(path)[:3]


def iter_shapes(shapes):
    """Recurse into group shapes so nested text/pictures are not missed."""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh.shapes)
        else:
            yield sh


def shape_text(sh):
    out = []
    if sh.has_text_frame:
        t = sh.text_frame.text.strip()
        if t:
            out.append(t)
    if sh.has_table:
        for row in sh.table.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                out.append(line)
    return out


def main():
    files = sorted(glob.glob(os.path.join(SRC, "M*_*.pptx")))
    files = [f for f in files if not os.path.basename(f).startswith("~$")]

    slides_out = open(os.path.join(OUT, "slides.jsonl"), "w")
    images_manifest = []
    summary = []
    seen_img_hash = {}

    for f in files:
        mod = module_id(f)
        prs = Presentation(f)
        deck_imgs = 0
        for si, slide in enumerate(prs.slides, start=1):
            texts, title = [], None
            # title placeholder if present
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()
            except Exception:
                pass
            img_idx = 0
            for sh in iter_shapes(slide.shapes):
                texts.extend(shape_text(sh))
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_idx += 1
                    deck_imgs += 1
                    try:
                        blob = sh.image.blob
                        ext = sh.image.ext
                        digest = hashlib.sha256(blob).hexdigest()
                        w_in = round(Emu(sh.width).inches, 2) if sh.width else 0
                        h_in = round(Emu(sh.height).inches, 2) if sh.height else 0
                        fname = f"{mod}_s{si:03d}_i{img_idx}.{ext}"
                        fpath = os.path.join(IMG_DIR, fname)
                        with open(fpath, "wb") as imf:
                            imf.write(blob)
                        images_manifest.append({
                            "file": fname,
                            "module": mod,
                            "deck": os.path.basename(f),
                            "slide_no": si,
                            "img_idx": img_idx,
                            "ext": ext,
                            "w_in": w_in,
                            "h_in": h_in,
                            "area_in2": round(w_in * h_in, 1),
                            "bytes": len(blob),
                            "sha256": digest[:16],
                            "dup_of": seen_img_hash.get(digest),
                        })
                        seen_img_hash.setdefault(digest, fname)
                    except Exception as e:
                        images_manifest.append({"file": None, "module": mod,
                                                "slide_no": si, "error": str(e)})
            # speaker notes
            notes = ""
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
            text_joined = "\n".join(texts).strip()
            if not title and texts:
                title = texts[0][:120]
            slides_out.write(json.dumps({
                "deck": os.path.basename(f),
                "module": mod,
                "slide_no": si,
                "title": title or "",
                "text": text_joined,
                "notes": notes,
            }, ensure_ascii=False) + "\n")
        summary.append({"module": mod, "deck": os.path.basename(f),
                        "slides": len(prs.slides), "images": deck_imgs})
        print(f"{mod} {os.path.basename(f):50s} slides={len(prs.slides):3d} imgs={deck_imgs:3d}")

    slides_out.close()
    json.dump(images_manifest, open(os.path.join(OUT, "images.json"), "w"), indent=2)
    json.dump(summary, open(os.path.join(OUT, "summary.json"), "w"), indent=2)
    uniq = len(set(i.get("sha256") for i in images_manifest if i.get("sha256")))
    print(f"\nWrote {len(images_manifest)} image records ({uniq} unique) to {IMG_DIR}")
    print(f"slides.jsonl, images.json, summary.json in {OUT}")


if __name__ == "__main__":
    main()
