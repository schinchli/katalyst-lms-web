#!/usr/bin/env python3
"""
Lean diagram shortlist (NO vision): score each slide of the target modules using
already-available signals, so we only thumbnail-confirm likely diagram slides.

Signals per slide (from PPTX, cheap):
  - biggest embedded raster image area (in^2)
  - embedded image count
  - connector count + group-shape count + autoshape count (shape-built diagrams)
  - total text length (diagrams have labels, not paragraphs)
Output: _extracted/shortlist.json  {module: [ {slide_no, score, signals...} ]}
"""
import glob, os, re, json
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = "/Users/schinchli/Documents/Training Material /Security Engineering on AWS"
OUT = os.path.join(SRC, "_extracted", "shortlist.json")
TARGET = {"M01", "M04", "M05", "M06", "M07", "M08"}


def walk(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(sh.shapes)


def score_slide(slide):
    max_area = 0.0
    imgs = connectors = groups = autoshapes = 0
    textlen = 0
    for sh in walk(slide.shapes):
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            imgs += 1
            try:
                a = Emu(sh.width).inches * Emu(sh.height).inches
                max_area = max(max_area, a)
            except Exception:
                pass
        elif st == MSO_SHAPE_TYPE.GROUP:
            groups += 1
        elif st in (MSO_SHAPE_TYPE.LINE,) or 'CONNECTOR' in str(st):
            connectors += 1
        elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            autoshapes += 1
        if getattr(sh, "has_text_frame", False):
            textlen += len(sh.text_frame.text)
    score = 0
    if max_area >= 6: score += 2
    if imgs >= 4: score += 2
    if connectors >= 2: score += 2
    if groups >= 1: score += 1
    if autoshapes >= 5 and textlen < 500: score += 1
    return score, dict(max_area=round(max_area, 1), imgs=imgs, connectors=connectors,
                       groups=groups, autoshapes=autoshapes, textlen=textlen)


def main():
    out = {}
    files = sorted(glob.glob(os.path.join(SRC, "M*_*.pptx")))
    files = [f for f in files if not os.path.basename(f).startswith("~$")]
    for f in files:
        mod = re.match(r"(M\d{2})_", os.path.basename(f)).group(1)
        if mod not in TARGET:
            continue
        prs = Presentation(f)
        cands = []
        for i, slide in enumerate(prs.slides, 1):
            sc, sig = score_slide(slide)
            if sc >= 2:
                cands.append({"slide_no": i, "score": sc, **sig})
        cands.sort(key=lambda c: -c["score"])
        out[mod] = cands
        print(f"{mod}: {len(cands)} candidate slides (of {len(prs.slides)})")
    json.dump(out, open(OUT, "w"), indent=2)
    tot = sum(len(v) for v in out.values())
    print(f"\nTotal candidates: {tot} -> {OUT}")


if __name__ == "__main__":
    main()
